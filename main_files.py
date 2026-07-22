from mitmproxy import http
from mitmproxy import ctx
import logging
import base64
import json
import os  # for database path
from datetime import datetime  # for recording capture time
import time
import io
import subprocess
import tempfile
import shutil
import xml.etree.ElementTree as ET
import zipfile
import re
from urllib.parse import urlparse, parse_qs
from dotenv import load_dotenv
import asyncio
from functools import partial

# SQLAlchemy core components
from sqlalchemy import create_engine, Column, Integer, String, DateTime, Text, LargeBinary
from sqlalchemy.orm import declarative_base, sessionmaker
from sqlalchemy import func  # import func for DB timestamp

# PDF parsing
try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

# Office document parsing
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# LibreOffice CLI (legacy format -> txt) — also check portable in project dir
SOFFICE_PORTABLE = os.path.join(os.path.dirname(__file__), 'libreoffice', 'program', 'soffice.exe')
HAS_SOFFICE = shutil.which('soffice') is not None or os.path.exists(SOFFICE_PORTABLE)

try:
    import xlrd
    HAS_XLRD = True
except ImportError:
    HAS_XLRD = False

def _soffice_cmd():
    if shutil.which('soffice'):
        return 'soffice'
    if os.path.exists(SOFFICE_PORTABLE):
        return SOFFICE_PORTABLE
    return None

# Local AI evaluation logic
from localAI_Reviw_files import my, my_image, my_classify_site, detect_and_mask

async def ai_call(func, *args, **kwargs):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(func, **kwargs), *args)

# ==================== 1. Database Configuration & Schema ====================
DB_PATH = os.path.join(os.path.abspath(os.path.dirname(__file__)), 'project.db')
ENGINE_URL = f'sqlite:///{DB_PATH}'
engine = create_engine(ENGINE_URL, echo=False) # echo=False to keep console clean
Base = declarative_base()

# Table 1: Admin users
class AdminUser(Base):
    __tablename__ = 'admin_users'
    
    id = Column(Integer, primary_key=True, autoincrement=True)
    username = Column(String(50), unique=True, nullable=False)
    password = Column(String(100), nullable=False)
    role = Column(String(20), server_default="user")
    created_at = Column(DateTime, server_default=func.now())
    department = Column(String(20), nullable=False)

# Table 2: Mitmproxy audit logs
class AuditLog(Base):
    __tablename__ = 'audit_logs'
    
    id = Column(Integer, primary_key=True, autoincrement=True)
    userid = Column(String(50))
    user_ip = Column(String(50))
    user_message = Column(Text)       # user chat content / file description
    review_status = Column(String(20)) # approved, warning, block
    review_reason = Column(Text)      # block reason
    raw_ai_json = Column(Text)        # full AI evaluation JSON
    captured_at = Column(DateTime, default=datetime.now)

# Table 3: File type registry (self-learning)
class FileTypeRegistry(Base):
    __tablename__ = 'file_type_registry'
    extension = Column(String(20), primary_key=True)
    is_document = Column(Integer, default=0)
    read_library = Column(String(100), default="")
    extraction_note = Column(Text, default="")
    discovered_at = Column(DateTime, default=datetime.now)

# Table 4: File records
class FileRecord(Base):
    __tablename__ = 'file_records'
    id = Column(Integer, primary_key=True, autoincrement=True)
    userid = Column(String(50))
    user_ip = Column(String(50))
    file_name = Column(String(255))
    file_type = Column(String(20))      # image / text / pdf / ooxml / unknown
    mime_type = Column(String(100))
    file_size = Column(Integer)
    file_data = Column(LargeBinary)
    extracted_text = Column(Text)
    review_status = Column(String(20))
    review_reason = Column(Text)
    raw_ai_json = Column(Text)
    captured_at = Column(DateTime, default=datetime.now)

# Table 5: Site registry
class SiteRegistry(Base):
    __tablename__ = 'site_registry'
    domain = Column(String(255), primary_key=True)
    is_ai = Column(Integer, default=0)          # 0=unknown 1=yes 2=no
    is_authorized = Column(Integer, default=0)  # 0=pending 1=yes 2=denied
    classification_reason = Column(Text, default="")
    search_summary = Column(Text, default="")
    tavily_raw = Column(Text, default="")
    discovered_at = Column(DateTime, default=datetime.now)
    reviewed_by = Column(String(50), default="")
    reviewed_at = Column(DateTime)

# Auto-create database & tables
Base.metadata.create_all(engine)
SessionLocal = sessionmaker(bind=engine)

# Seed admin account (if not exists)
db_session = SessionLocal()
if not db_session.query(AdminUser).filter_by(username='admin').first():
    admin = AdminUser(id=0, username='admin', password='admin', role='admin', department='IT') # test password
    db_session.add(admin)
    db_session.commit()
db_session.close()
# ==================================================================

# Logging config
logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)
log_format = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')

# ---- Handler 1: Console ----
c_handler = logging.StreamHandler()
c_handler.setFormatter(log_format)
c_handler.setLevel(logging.INFO)

# ---- Handler 3: WARNING/ERROR file ----
warn_handler = logging.FileHandler("warning.log", encoding="utf-8")
warn_handler.setFormatter(log_format)
warn_handler.setLevel(logging.WARNING)

logger.addHandler(c_handler)
logger.addHandler(warn_handler)


# ==================== 2. Global Connection State ====================
CONNECTION_STATES = {}
IP_AUTH_CACHE = {} 
IP_CACHE_TTL = 300  # cache TTL 5 min
# (SITE_CACHE removed — background classification replaces sync blocking)

class Auth():
    def __init__(self):
        logger.info("Proxy Authentication enabled via Strict HTTP CONNECT & Request hooks.")

    def _authenticate(self, flow: http.HTTPFlow) -> tuple[bool, str, str]:
        """Internal: extract & validate credentials from request headers"""
        proxy_auth = flow.request.headers.get("Proxy-Authorization")
        
        if not proxy_auth or not proxy_auth.startswith("Basic "):
            return False, "N/A", "user"
        
        try:
            b64_token = proxy_auth.split(" ", 1)[1]
            decoded = base64.b64decode(b64_token).decode("utf-8", errors="ignore")
            input_username, input_password = decoded.split(":", 1)
        except Exception as e:
            logger.error(f"Failed to parse proxy credentials: {e}")
            return False, "N/A", "user"
        
        session = SessionLocal()
        db_user = session.query(AdminUser).filter_by(username=input_username).first()
        session.close()

        if db_user and db_user.password == input_password:
            return True, input_username, (db_user.role if db_user.role else "user")
        return False, input_username, "user"

    def http_connect(self, flow: http.HTTPFlow) -> None:
        """Core: intercept HTTPS CONNECT requests"""
        success, username, role = self._authenticate(flow)
        client_peer = flow.client_conn.peername
        
        if success:
            CONNECTION_STATES[flow.client_conn.id] = {
                "auth_success": True,
                "userid": username,
                "user_ip": str(client_peer),
                "user_role": role
            }
            if client_peer:
                IP_AUTH_CACHE[client_peer[0]] = {
                    "userid": username,
                    "user_role": role,
                    "timestamp": time.time()
                }
            logger.info(f"[HTTPS tunnel allowed] user: {username} | role: {role} | IP: {client_peer}")
        else:
            if username != "N/A":
                logger.warning(f"[HTTPS auth denied] rejected login attempt -> account: {username}")
            
            flow.response = http.Response.make(
                407, b"Proxy Authentication Required",
                {
                    b"Proxy-Authenticate": b'Basic realm="proxy auth required"',
                    b"Connection": b"close"
                }
            )

    def request(self, flow: http.HTTPFlow) -> None:
        """Second layer: handle individual requests"""
        if flow.response and flow.response.status_code == 407:
            return

        client_peer = flow.client_conn.peername
        client_ip = client_peer[0] if client_peer else None

        # Logic A: HTTPS tunnel traffic, inherit metadata from connection dict
        conn_state = CONNECTION_STATES.get(flow.client_conn.id, {})
        if conn_state.get("auth_success", False):
            flow.metadata["userid"] = conn_state.get("userid", "N/A")
            flow.metadata["user_ip"] = conn_state.get("user_ip", "Unknown")
            flow.metadata["user_role"] = conn_state.get("user_role", "user")
            return

        # Logic B: Plaintext request with header credentials, validate & pass
        success, username, role = self._authenticate(flow)
        if success:
            flow.metadata["userid"] = username
            flow.metadata["user_ip"] = str(client_peer)
            flow.metadata["user_role"] = role
            if client_ip:
                IP_AUTH_CACHE[client_ip] = {
                    "userid": username,
                    "user_role": role,
                    "timestamp": time.time()
                }
            return
        elif username != "N/A":
            logger.warning(f"[HTTP auth denied] rejected plaintext request -> account: {username}")

        # Logic C: IP-level sliding window reuse
        if client_ip and client_ip in IP_AUTH_CACHE:
            cache_info = IP_AUTH_CACHE[client_ip]
            if time.time() - cache_info["timestamp"] < IP_CACHE_TTL:
                cache_info["timestamp"] = time.time()
                flow.metadata["userid"] = cache_info["userid"]
                flow.metadata["user_ip"] = str(client_peer)
                flow.metadata["user_role"] = cache_info["user_role"]
                return

        # Logic D: Deny unauthorized access
        logger.error(f"[UNAUTHORIZED] Detected unauthorized request! Source IP: {client_peer}")
        flow.response = http.Response.make(
            407, b"Proxy Authentication Required",
            {
                b"Proxy-Authenticate": b'Basic realm="Mitmproxy Admin Required"',
                b"Connection": b"close"
            }
        )

    def client_disconnected(self, client) -> None:
        """Memory cleanup"""
        CONNECTION_STATES.pop(client.id, None)


TARGET_PATHS = ["/backend-api/f/conversation", "/backend-anon/f/conversation"]

def extract_text_from_pdf(pdf_bytes):
    """Extract text from raw PDF bytes"""
    if not HAS_PYPDF:
        logger.error("pypdf is not installed. PDF text extraction is skipped.")
        return ""
    try:
        pdf_file = io.BytesIO(pdf_bytes)
        reader = pypdf.PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text.strip()
    except Exception as e:
        logger.error(f"Failed to extract PDF text: {e}")
        return ""


MIME_EXT_MAP = {
    'application/msword': 'doc',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.ms-excel': 'xls',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'application/vnd.ms-powerpoint': 'ppt',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'application/vnd.oasis.opendocument.text': 'odt',
    'application/vnd.oasis.opendocument.spreadsheet': 'ods',
    'application/vnd.oasis.opendocument.presentation': 'odp',
    'application/rtf': 'rtf',
    'text/plain': 'txt',
    'text/csv': 'csv',
    'application/pdf': 'pdf',
}

def get_extension(flow):
    path = flow.request.path
    qs = parse_qs(urlparse(flow.request.url).query)
    filename = qs.get('filename', [None])[0]
    if filename:
        ext = os.path.splitext(filename)[1].lower().lstrip('.')
        if ext:
            return ext
    ext = os.path.splitext(path.split('/')[-1])[1].lower().lstrip('.')
    if ext:
        return ext
    cd = flow.request.headers.get("Content-Disposition", "")
    m = re.search(r'filename=["\']?([^"\';\n]+)["\']?', cd)
    if m:
        ext = os.path.splitext(m.group(1))[1].lower().lstrip('.')
        return ext
    ct = flow.request.headers.get("content-type", "").lower().split(';')[0].strip()
    if ct in MIME_EXT_MAP:
        return MIME_EXT_MAP[ct]
    return ""


def get_filename(flow, ext=""):
    qs = parse_qs(urlparse(flow.request.url).query)
    filename = qs.get('filename', [None])[0]
    if filename:
        return filename
    path_fname = flow.request.path.split('/')[-1]
    if path_fname and '.' in path_fname:
        return path_fname
    cd = flow.request.headers.get("Content-Disposition", "")
    m = re.search(r'filename=["\']?([^"\';\n]+)["\']?', cd)
    if m:
        return m.group(1)
    return f"unnamed.{ext}" if ext else "unnamed"

def try_decode_text(data):
    for enc in ['utf-8', 'utf-16']:
        try:
            text = data.decode(enc).strip()
            if text:
                return text
        except:
            continue
    return ""


def is_pdf(data):
    return data.startswith(b'%PDF')


def is_zip(data):
    return data[:4] == b'PK\x03\x04'


def extract_text_from_ooxml(data):
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = zf.namelist()
    except Exception:
        return ""
    if '[Content_Types].xml' not in names:
        return ""
    for extractor in (extract_text_from_docx, extract_text_from_pptx, extract_text_from_xlsx):
        text = extractor(data)
        if text:
            return text
    return ""


def extract_text_from_docx(data):
    if not HAS_DOCX:
        logger.error("python-docx not installed, cannot extract .docx text")
        return ""
    try:
        doc = docx.Document(io.BytesIO(data))
        texts = []
        for p in doc.paragraphs:
            if p.text:
                texts.append(p.text)
        for t in doc.tables:
            for r in t.rows:
                for c in r.cells:
                    for p in c.paragraphs:
                        if p.text:
                            texts.append(p.text)
        return '\n'.join(texts).strip()
    except Exception as e:
        logger.error(f"python-docx error: {e}")
        return ""


def extract_text_from_pptx(data):
    if not HAS_PPTX:
        logger.error("python-pptx not installed, cannot extract .pptx text")
        return ""
    try:
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text:
                            texts.append(p.text)
                if shape.has_table:
                    for r in shape.table.rows:
                        for c in r.cells:
                            if c.text:
                                texts.append(c.text)
        return '\n'.join(texts).strip()
    except Exception as e:
        logger.error(f"python-pptx error: {e}")
        return ""


def extract_text_from_xlsx(data):
    if not HAS_OPENPYXL:
        logger.error("openpyxl not installed, cannot extract .xlsx text")
        return ""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        texts = []
        for name in wb.sheetnames:
            ws = wb[name]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        texts.append(str(cell.value))
        wb.close()
        return '\n'.join(texts).strip()
    except Exception as e:
        logger.error(f"openpyxl error: {e}")
        return ""


def extract_text_from_xls(data):
    if not HAS_XLRD:
        return ""
    try:
        wb = xlrd.open_workbook(file_contents=data)
        texts = []
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            for row in range(ws.nrows):
                for col in range(ws.ncols):
                    cell = ws.cell(row, col)
                    if cell.value:
                        texts.append(str(cell.value))
        return '\n'.join(texts).strip()
    except Exception as e:
        logger.error(f"xlrd error: {e}")
        return ""


def extract_text_with_libreoffice(data, ext):
    if not HAS_SOFFICE:
        return ""
    tmpdir = None
    try:
        tmpdir = tempfile.mkdtemp(prefix='office_conv_')
        in_path = os.path.join(tmpdir, f'input.{ext}' if ext else 'input.bin')
        with open(in_path, 'wb') as f:
            f.write(data)
        subprocess.run(
            [_soffice_cmd(), '--headless', '--convert-to', 'txt:Text', '--outdir', tmpdir, in_path],
            capture_output=True, timeout=30
        )
        out_path = os.path.splitext(in_path)[0] + '.txt'
        if os.path.exists(out_path):
            with open(out_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read().strip()
        return ""
    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
        return ""
    except Exception as e:
        logger.error(f"LibreOffice conversion error: {e}")
        return ""
    finally:
        if tmpdir:
            shutil.rmtree(tmpdir, ignore_errors=True)


OOXML_MIME_MAP = {
    '.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
    '.gif': 'image/gif', '.bmp': 'image/bmp', '.webp': 'image/webp',
    '.tiff': 'image/tiff', '.tif': 'image/tiff',
}


def extract_media_from_ooxml(data):
    images = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            for name in zf.namelist():
                if name.startswith(('word/media/', 'ppt/media/', 'xl/media/')):
                    img_data = zf.read(name)
                    ext = os.path.splitext(name)[1].lower()
                    mime = OOXML_MIME_MAP.get(ext, 'image/png')
                    images.append((img_data, mime, name))
    except Exception as e:
        logger.error(f"OOXML media extraction error: {e}")
    return images


def extract_strings(data, min_len=4):
    result = []
    buf = []
    for b in data:
        if 32 <= b < 127:
            buf.append(chr(b))
        else:
            if len(buf) >= min_len:
                result.append(''.join(buf))
            buf = []
    if len(buf) >= min_len:
        result.append(''.join(buf))
    return '\n'.join(result)


def classify_extension(ext):
    if not ext:
        return False
    session = SessionLocal()
    try:
        existing = session.query(FileTypeRegistry).filter_by(extension=ext).first()
        if existing:
            return existing.is_document == 1
        from localAI_Reviw_files import my_classify
        result = my_classify(ext)
        entry = FileTypeRegistry(
            extension=ext,
            is_document=1 if result.get('is_document') else 2,
            read_library=result.get('read_library', ''),
            extraction_note=result.get('note', ''),
            discovered_at=datetime.now()
        )
        session.add(entry)
        session.commit()
        return result.get('is_document', False)
    except Exception as e:
        session.rollback()
        logger.error(f"Extension classification error: {e}")
        return False
    finally:
        session.close()


class AI_checkpoint():
    def __init__(self):
        logger.info("Audit Service Started (Text + Images + Documents)")

    async def request(self, flow: http.HTTPFlow) -> None:
        # ==================== Intercept Text Messages ====================
        if flow.request.method == "POST" and any(path in flow.request.path for path in TARGET_PATHS):
            if not flow.request.text:
                return
            
            try:
                userid = flow.metadata.get("userid", "N/A")
                user_ip = flow.metadata.get("user_ip", "Unknown")

                data = json.loads(flow.request.text)
                messages = data.get("messages", [])
                for message in messages:
                    author = message.get("author", {})
                    if author.get("role") == "user":
                        content = message.get("content", {})
                        if content.get("content_type") == "text":
                            parts = content.get("parts", [])
                            text_parts = [p for p in parts if isinstance(p, str)]
                            user_message = "".join(text_parts).strip()
                            
                            if not user_message:
                                continue
                            
                            # Send masked content to server
                            detect_result = await ai_call(detect_and_mask, user_message)
                            masked_message = detect_result.get("masked", user_message)
                            if detect_result.get("has_sensitive"):
                                content["parts"] = [masked_message]
                                flow.request.text = json.dumps(data, ensure_ascii=False)
                            
                            msg_review_result = await ai_call(my, masked_message, skip_mask=True)
                            review_status = msg_review_result["status"]
                            review_reason = msg_review_result["reason"]
                            logger.warning(msg_review_result)
                            logger.warning(f"UserID: {userid} IP: {user_ip} | Msg: {user_message} | Status: {review_status} | Reason: {review_reason}")

                            session = SessionLocal()
                            try:
                                log_entry = AuditLog(
                                    userid=userid,
                                    user_ip=user_ip,
                                    user_message=user_message,
                                    review_status=review_status,
                                    review_reason=review_reason,
                                    raw_ai_json=json.dumps(msg_review_result, ensure_ascii=False)
                                )
                                session.add(log_entry)
                                session.commit()
                            except Exception as db_err:
                                session.rollback()
                                logger.error(f"DB write failed: {db_err}")
                            finally:
                                session.close()
                            
                            if review_status == "block":
                                flow.response = http.Response.make(
                                    403,
                                    b"Request blocked by local policy. Please do not send aggressive content.",
                                    {b"Content-Type": b"text/plain; charset=utf-8"}
                                )
                                return
            except json.JSONDecodeError:
                pass
            except Exception as e:
                logger.error(f"Error processing text request: {e}")

        # ==================== Intercept Images/Docs/Files (fallback extraction pipeline) ====================
        elif flow.request.method == "PUT" and "oaiusercontent.com" in flow.request.pretty_host and "/files/" in flow.request.path and "/raw" in flow.request.path:
            content_type = flow.request.headers.get("content-type", "").lower()
            userid = flow.metadata.get("userid", "N/A")
            user_ip = flow.metadata.get("user_ip", "Unknown")
            file_bytes = flow.request.content

            if not file_bytes:
                return

            ext = get_extension(flow)
            file_name = get_filename(flow, ext)
            file_type = "unknown"
            extracted_text = None
            review_status = "approved"
            review_reason = ""
            raw_ai_json = "{}"
            file_description = f"File (MimeType: {content_type}, Ext: .{ext}, Size: {len(file_bytes)} bytes)"

            # A. Image -> Vision AI review
            if content_type.startswith("image/"):
                file_type = "image"
                logger.info(f"[*] Evaluating uploaded image via local AI (size: {len(file_bytes)} bytes)...")
                try:
                    img_review_result = await ai_call(my_image, file_bytes, content_type)
                    review_status = img_review_result.get("status", "block")
                    review_reason = img_review_result.get("reason", "No reason provided")
                    raw_ai_json = json.dumps(img_review_result, ensure_ascii=False)
                    file_description = f"[Image Upload] MimeType: {content_type}, Size: {len(file_bytes)} bytes"
                except Exception as e:
                    logger.error(f"Local AI image review error: {e}")
                    review_status = "block"
                    review_reason = f"AI Image review error: {str(e)}"

            # B-D: Text / PDF / OOXML / Unknown -- magic-byte fallback extraction
            else:
                # == B. Try text decode first ==
                text = try_decode_text(file_bytes)
                if text:
                    file_type = "text"
                    extracted_text = text
                    logger.info(f"[*] Evaluating uploaded text file via local AI (size: {len(file_bytes)} bytes)...")
                    try:
                        detect_sensitive = await ai_call(detect_and_mask, text)
                        if detect_sensitive.get("has_sensitive"):
                            review_status = "block"
                            review_reason = "File contains sensitive content, cannot be masked"
                            raw_ai_json = json.dumps(detect_sensitive, ensure_ascii=False)
                            file_description = f"[Text File Upload - Blocked] Ext: .{ext}, Size: {len(file_bytes)} bytes"
                        else:
                            msg_review_result = await ai_call(my, text, skip_mask=True)
                            review_status = msg_review_result.get("status", "block")
                            review_reason = msg_review_result.get("reason", "No reason provided")
                            raw_ai_json = json.dumps(msg_review_result, ensure_ascii=False)
                            file_description = f"[Text File Upload] Ext: .{ext}, Size: {len(file_bytes)} bytes"
                    except Exception as e:
                        logger.error(f"Local AI text review error: {e}")
                        review_status = "block"
                        review_reason = f"AI Text file review error: {str(e)}"

                # == C. PDF ==
                elif is_pdf(file_bytes):
                    file_type = "pdf"
                    logger.info(f"[*] Evaluating uploaded PDF via local AI (size: {len(file_bytes)} bytes)...")
                    try:
                        if not HAS_PYPDF:
                            review_status = "block"
                            review_reason = "PDF uploaded but pypdf is not installed on proxy server"
                        else:
                            pdf_text = await ai_call(extract_text_from_pdf, file_bytes)
                            if pdf_text:
                                extracted_text = pdf_text
                                logger.info(f"[*] PDF extracted {len(pdf_text)} chars, sending to local AI...")
                                detect_sensitive = await ai_call(detect_and_mask, pdf_text)
                                if detect_sensitive.get("has_sensitive"):
                                    review_status = "block"
                                    review_reason = "PDF contains sensitive content, cannot be masked"
                                    raw_ai_json = json.dumps(detect_sensitive, ensure_ascii=False)
                                else:
                                    msg_review_result = await ai_call(my, pdf_text, skip_mask=True)
                                    review_status = msg_review_result.get("status", "block")
                                    review_reason = msg_review_result.get("reason", "No reason provided")
                                    raw_ai_json = json.dumps(msg_review_result, ensure_ascii=False)
                            else:
                                logger.warning("[!] PDF text extraction returned empty (likely image-based PDF)")
                                review_status = "approved"
                                review_reason = "PDF has no extractable text content"
                            file_description = f"[PDF Upload] Size: {len(file_bytes)} bytes, Extracted: {len(pdf_text)} chars"
                    except Exception as e:
                        logger.error(f"Local AI PDF review error: {e}")
                        review_status = "block"
                        review_reason = f"AI PDF review error: {str(e)}"

                # == D. OOXML (docx/pptx/xlsx) ==
                elif is_zip(file_bytes):
                    file_type = "ooxml"
                    logger.info(f"[*] Parsing OOXML document (.docx/.pptx/.xlsx)...")
                    try:
                        # Review embedded images first (word/media/ etc.)
                        ooxml_images = await ai_call(extract_media_from_ooxml, file_bytes)
                        if ooxml_images:
                            logger.info(f"[*] OOXML has {len(ooxml_images)} embedded images, reviewing...")
                            for img_data, img_mime, img_name in ooxml_images:
                                result = await ai_call(my_image, img_data, img_mime)
                                if result.get("status") == "block":
                                    review_status = "block"
                                    review_reason = f"Embedded image blocked: {img_name} - {result.get('reason', '')}"
                                    raw_ai_json = json.dumps(result, ensure_ascii=False)
                                    logger.warning(f"[!] OOXML embedded image blocked: {img_name}")
                                    break
                        # Then review text
                        if review_status != "block":
                            ooxml_text = await ai_call(extract_text_from_ooxml, file_bytes)
                            if ooxml_text:
                                extracted_text = ooxml_text
                                logger.info(f"[*] OOXML extracted {len(ooxml_text)} chars, sending to local AI...")
                                detect_sensitive = await ai_call(detect_and_mask, ooxml_text)
                                if detect_sensitive.get("has_sensitive"):
                                    review_status = "block"
                                    review_reason = "OOXML document contains sensitive content, cannot be masked"
                                    raw_ai_json = json.dumps(detect_sensitive, ensure_ascii=False)
                                else:
                                    msg_review_result = await ai_call(my, ooxml_text, skip_mask=True)
                                    review_status = msg_review_result.get("status", "block")
                                    review_reason = msg_review_result.get("reason", "No reason provided")
                                    raw_ai_json = json.dumps(msg_review_result, ensure_ascii=False)
                            else:
                                logger.warning("[!] OOXML document yielded no text")
                                if review_status != "block":
                                    review_status = "approved"
                                    review_reason = "OOXML document has no extractable text content"
                        file_description = f"[OOXML Upload] Ext: .{ext}, Size: {len(file_bytes)} bytes, Images: {len(ooxml_images)}"
                    except Exception as e:
                        logger.error(f"OOXML parse error: {e}")
                        review_status = "block"
                        review_reason = f"OOXML parse error: {str(e)}"

                # == E. Fallback: unknown binary (AI learning + LibreOffice + strings) ==
                else:
                    file_type = "unknown"
                    is_doc = classify_extension(ext)
                    if is_doc:
                        lo_text = ""
                        # 0) xlrd for .xls files
                        if ext == 'xls' and HAS_XLRD:
                            xls_text = await ai_call(extract_text_from_xls, file_bytes)
                            if xls_text:
                                extracted_text = xls_text
                                logger.info(f"[*] xlrd converted .xls to text ({len(xls_text)} chars), sending to AI...")
                                detect_sensitive = await ai_call(detect_and_mask, xls_text)
                                if detect_sensitive.get("has_sensitive"):
                                    review_status = "block"
                                    review_reason = "Document contains sensitive content, cannot be masked"
                                    raw_ai_json = json.dumps(detect_sensitive, ensure_ascii=False)
                                else:
                                    msg_review_result = await ai_call(my, xls_text, skip_mask=True)
                                    review_status = msg_review_result.get("status", "block")
                                    review_reason = msg_review_result.get("reason", "No reason provided")
                                    raw_ai_json = json.dumps(msg_review_result, ensure_ascii=False)
                        if not extracted_text:
                            lo_text = await ai_call(extract_text_with_libreoffice, file_bytes, ext)
                        if lo_text:
                            extracted_text = lo_text
                            logger.info(f"[*] LibreOffice converted .{ext} to text ({len(lo_text)} chars), sending to AI...")
                            # Log to file_type_registry for admin review
                            try:
                                sess = SessionLocal()
                                entry = sess.query(FileTypeRegistry).filter_by(extension=ext).first()
                                if entry and entry.read_library != 'LibreOffice CLI':
                                    entry.read_library = 'LibreOffice CLI'
                                    entry.extraction_note = 'LibreOffice headless --convert-to txt:Text'
                                    sess.commit()
                                sess.close()
                            except Exception:
                                pass
                            detect_sensitive = await ai_call(detect_and_mask, lo_text)
                            if detect_sensitive.get("has_sensitive"):
                                review_status = "block"
                                review_reason = "Document contains sensitive content, cannot be masked"
                                raw_ai_json = json.dumps(detect_sensitive, ensure_ascii=False)
                            else:
                                msg_review_result = await ai_call(my, lo_text, skip_mask=True)
                                review_status = msg_review_result.get("status", "block")
                                review_reason = msg_review_result.get("reason", "No reason provided")
                                raw_ai_json = json.dumps(msg_review_result, ensure_ascii=False)
                        else:
                            logger.warning(f"[!] LibreOffice .{ext} conversion failed, falling back to strings...")
                            strings_text = await ai_call(extract_strings, file_bytes)
                            if strings_text:
                                extracted_text = strings_text
                                detect_sensitive = await ai_call(detect_and_mask, strings_text)
                                if detect_sensitive.get("has_sensitive"):
                                    review_status = "block"
                                    review_reason = "Unknown document contains sensitive content, cannot be masked"
                                    raw_ai_json = json.dumps(detect_sensitive, ensure_ascii=False)
                                else:
                                    msg_review_result = await ai_call(my, strings_text, skip_mask=True)
                                    review_status = msg_review_result.get("status", "block")
                                    review_reason = msg_review_result.get("reason", "No reason provided")
                                    raw_ai_json = json.dumps(msg_review_result, ensure_ascii=False)
                            else:
                                review_status = "approved"
                                review_reason = f"Unsupported document format .{ext}, no text extractable (see file_type_registry in DB)"
                    else:
                        review_status = "approved"
                        review_reason = f"Non-document format .{ext} (is_document={is_doc}), skipped"
                    file_description = f"[Unknown Format Upload] Ext: .{ext}, Size: {len(file_bytes)} bytes"

            # Print/audit result
            logger.warning(f"File Review: UserID: {userid} IP: {user_ip} | Status: {review_status} | Reason: {review_reason}")

            # Save to DB audit table
            session = SessionLocal()
            try:
                log_entry = AuditLog(
                    userid=userid,
                    user_ip=user_ip,
                    user_message=file_description,
                    review_status=review_status,
                    review_reason=review_reason,
                    raw_ai_json=raw_ai_json
                )
                session.add(log_entry)
                file_record = FileRecord(
                    userid=userid, user_ip=user_ip,
                    file_name=file_name, file_type=file_type,
                    mime_type=content_type, file_size=len(file_bytes),
                    file_data=file_bytes,
                    extracted_text=extracted_text,
                    review_status=review_status,
                    review_reason=review_reason, raw_ai_json=raw_ai_json
                )
                session.add(file_record)
                session.commit()
            except Exception as db_err:
                session.rollback()
                logger.error(f"DB record failed: {db_err}")
            finally:
                session.close()

            # Block if needed
            if review_status == "block":
                flow.response = http.Response.make(
                    403,
                    f"Blocked: File failed security review. Reason: {review_reason}".encode("utf-8"),
                    {b"Content-Type": b"text/plain; charset=utf-8"}
                )
                return


# ==================== 3. AI Site Interceptor ====================

def classify_site(host):
    from ddgs import DDGS
    try:
        raw_results = []
        with DDGS(timeout=10) as ddgs:
            try:
                raw_results = list(ddgs.text(f"{host} website definition overview", max_results=5))
            except Exception:
                raw_results = []
        search_results_str = ""
        if raw_results:
            for idx, res in enumerate(raw_results, 1):
                title = res.get('title', 'No Title')
                snippet = res.get('body', 'No Snippet')
                search_results_str += f"[{idx}] Title: {title}\nSnippet: {snippet}\n\n"
        else:
            search_results_str = "No internet search results available."
        context = (
            f"Task: Please review the following search results about the website '{host}'.\n"
            f"Determine if this website is related to AI chatbots.\n"
            f"If it is related to AI chatbots, reply ONLY 'True'. If not, reply 'False'.\n\n"
            f"Search Results:\n{search_results_str}"
        )
        result = my_classify_site(host, context)
        logger.debug(f"[SITE DEBUG] AI verdict for {host}: {json.dumps(result, ensure_ascii=False)}")
        return result.get("is_ai", False), result.get("reason", ""), json.dumps(raw_results, ensure_ascii=False)
    except Exception as e:
        logger.error(f"Site classify error for {host}: {e}")
        return False, "", ""


_classifying_hosts = set()

class AISiteChecker():
    def __init__(self):
        logger.info("AI Site Checker started — non-AI sites pass, unknown auto-classify, unapproved AI blocked")

    async def request(self, flow: http.HTTPFlow) -> None:
        if flow.response:
            return
        host = flow.request.pretty_host
        if not host:
            return
        userid = flow.metadata.get("userid", "N/A")
        user_ip = flow.metadata.get("user_ip", "Unknown")

        session = SessionLocal()
        try:
            entry = session.query(SiteRegistry).filter_by(domain=host).first()
        finally:
            session.close()

        if entry:
            # is_ai=0 classifying in bg -> pass
            if entry.is_ai == 1 and entry.is_authorized == 1:
                return
            if entry.is_ai == 2:
                return
            if entry.is_ai == 0:
                return
            logger.warning(f"[SITE BLOCK] Unauthorized AI site: {host} from {userid}@{user_ip}")
            flow.response = http.Response.make(
                403,
                f"Blocked: Unauthorized AI website ({host}). Contact admin.".encode("utf-8"),
                {b"Content-Type": b"text/plain; charset=utf-8"}
            )
            return

        # Not in DB -> insert (is_ai=0) -> pass -> background classify
        session = SessionLocal()
        try:
            session.add(SiteRegistry(
                domain=host, is_ai=0, is_authorized=0,
                discovered_at=datetime.now()
            ))
            session.commit()
        except Exception:
            session.rollback()
        finally:
            session.close()

        logger.info(f"[SITE ALLOW] New site (pending classify): {host} from {userid}@{user_ip}")

        if host not in _classifying_hosts:
            _classifying_hosts.add(host)
            asyncio.create_task(self._background_classify(host, userid, user_ip))

    async def _background_classify(self, host, userid, user_ip):
        try:
            logger.info(f"[SITE CLASSIFY] Background classification: {host}")
            is_ai, reason, tavily_raw = await ai_call(classify_site, host)
            session = SessionLocal()
            try:
                entry = session.query(SiteRegistry).filter_by(domain=host).first()
                if entry:
                    entry.is_ai = 1 if is_ai else 2
                    entry.classification_reason = reason
                    entry.tavily_raw = tavily_raw
                    session.commit()
                if is_ai:
                    logger.warning(f"[SITE CLASSIFY] {host} is AI site (pending admin approval)")
                else:
                    logger.info(f"[SITE CLASSIFY] {host} is non-AI site")
            except Exception as e:
                session.rollback()
                logger.error(f"SiteRegistry update error for {host}: {e}")
            finally:
                session.close()
        finally:
            _classifying_hosts.discard(host)


addons = [
    Auth(),
    AISiteChecker(),
    AI_checkpoint()
]
